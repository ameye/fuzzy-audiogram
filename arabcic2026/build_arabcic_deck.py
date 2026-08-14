#!/usr/bin/env python3
"""Build ArabCIC 2026 presentation deck from the fuzzy audiogram manuscript."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# ---------------------------------------------------------------- palette
NAVY   = RGBColor(0x0B, 0x25, 0x45)   # dark navy
NAVY2  = RGBColor(0x14, 0x38, 0x62)   # lighter navy (chips on dark)
TEAL   = RGBColor(0x1D, 0x78, 0x74)   # teal
GOLD   = RGBColor(0xE8, 0xB5, 0x4A)   # arab gold
LIGHT  = RGBColor(0xF6, 0xF8, 0xFA)   # light background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x1B, 0x2A, 0x3A)   # body text
MUTED  = RGBColor(0x5A, 0x6B, 0x7B)   # muted text
BORDER = RGBColor(0xD8, 0xE0, 0xE8)
SOFT_T = RGBColor(0xE4, 0xF0, 0xF0)   # soft teal card bg
SOFT_G = RGBColor(0xFB, 0xF3, 0xE0)   # soft gold card bg
CARD   = RGBColor(0xFF, 0xFF, 0xFF)

HDR = "Georgia"
BODY = "Calibri"

SLIDE_W, SLIDE_H = 13.333, 7.5
FIG = "/opt/data/fuzzy-audiogram/figures"

prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------- helpers
def add_slide(bg):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s

def rect(s, x, y, w, h, fill, line=None, rounded=False, radius=0.06, line_w=0.75):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def txt(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        wrap=True, m=Pt(4)):
    """paras: list of dicts: text|runs, size, bold, color, font, align, space_after, italic, line_spacing"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = m; tf.margin_right = m; tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", align)
        if p.get("space_after") is not None: para.space_after = Pt(p["space_after"])
        if p.get("space_before") is not None: para.space_before = Pt(p["space_before"])
        if p.get("line_spacing"): para.line_spacing = p["line_spacing"]
        runs = p.get("runs") or [{"text": p.get("text", ""), "size": p.get("size", 14),
                                  "bold": p.get("bold", False), "color": p.get("color", INK),
                                  "font": p.get("font", BODY), "italic": p.get("italic", False)}]
        for rr in runs:
            r = para.add_run(); r.text = rr.get("text", "")
            f = r.font
            f.size = Pt(rr.get("size", p.get("size", 14)))
            f.bold = rr.get("bold", p.get("bold", False))
            f.italic = rr.get("italic", p.get("italic", False))
            f.color.rgb = rr.get("color", p.get("color", INK))
            f.name = rr.get("font", p.get("font", BODY))
    return tb

def pic(s, path, x, y, w=None, h=None):
    im = Image.open(path); ar = im.size[0] / im.size[1]
    if w and not h: h = w / ar
    if h and not w: w = h * ar
    return s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h)), w, h

def chip(s, x, y, w, h, big, small, fill=NAVY2, big_color=GOLD, small_color=WHITE, big_size=30, small_size=10.5):
    c = rect(s, x, y, w, h, fill, rounded=True, radius=0.12)
    txt(s, x + 0.15, y + 0.04, w - 0.3, h * 0.52,
        [{"text": big, "size": big_size, "bold": True, "color": big_color,
          "font": HDR, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 0.15, y + h * 0.56, w - 0.3, h * 0.40,
        [{"text": small, "size": small_size, "color": small_color, "align": PP_ALIGN.CENTER}],
        anchor=MSO_ANCHOR.MIDDLE)

def waveform(s, x, y, w, color=GOLD, max_h=0.55, n=11):
    heights = [0.25, 0.5, 0.85, 1.15, 1.5, 1.75, 1.5, 1.15, 0.85, 0.5, 0.25]
    total_bars = n
    bw = 0.10; gap = 0.13
    seq_w = total_bars * (bw + gap)
    start = x + (w - seq_w) / 2
    for i in range(total_bars):
        hgt = max_h * heights[i % len(heights)]
        rect(s, start + i * (bw + gap), y - hgt, bw, hgt, color)

def title_bar(s, text, num=None, total=15):
    txt(s, 0.55, 0.30, 10.5, 0.75,
        [{"text": text, "size": 28, "bold": True, "color": NAVY, "font": HDR}])
    txt(s, 11.4, 0.42, 1.4, 0.4,
        [{"text": f"{num} / {total}", "size": 10, "color": MUTED, "align": PP_ALIGN.RIGHT}])
    # footer
    txt(s, 0.55, 7.08, 6.5, 0.3,
        [{"text": "ArabCIC 2026 · Jeddah · Fuzzy Audiometric Index (FAI)", "size": 9, "color": MUTED}])

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def card_header(s, x, y, w, h, label, fill=NAVY, color=WHITE, size=12.5):
    c = rect(s, x, y, w, h, fill, rounded=True, radius=0.10)
    txt(s, x + 0.1, y, w - 0.2, h,
        [{"text": label, "size": size, "bold": True, "color": color, "align": PP_ALIGN.CENTER}],
        anchor=MSO_ANCHOR.MIDDLE, m=Pt(2))

def make_table(s, x, y, w, data, col_w, row_h=0.42, header_fill=NAVY,
               body_size=11, header_size=11.5, highlight_rows=None):
    rows, cols = len(data), len(data[0])
    gf = s.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(row_h * rows))
    tbl = gf.table
    tbl.first_row = False; tbl.horz_banding = False
    for ci, cw in enumerate(col_w):
        tbl.columns[ci].width = Inches(cw)
    for ri in range(rows):
        tbl.rows[ri].height = Inches(row_h)
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Pt(6); cell.margin_right = Pt(6)
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                if highlight_rows and ri in highlight_rows:
                    cell.fill.solid(); cell.fill.fore_color.rgb = SOFT_G
                else:
                    cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xEF, 0xF4, 0xF8)
            tf = cell.text_frame; tf.word_wrap = True
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            r = para.add_run(); r.text = str(data[ri][ci])
            f = r.font
            f.size = Pt(header_size if ri == 0 else body_size)
            f.bold = (ri == 0)
            f.color.rgb = WHITE if ri == 0 else INK
            f.name = BODY
    return gf

# ================================================================ SLIDE 1 — TITLE (dark)
s = add_slide(NAVY)
waveform(s, 0, 7.28, 13.333, color=GOLD, max_h=0.55)
waveform(s, 0, 7.28, 13.333, color=TEAL, max_h=0.30)
txt(s, 0.8, 0.85, 11.73, 0.45,
    [{"text": "4TH ARAB COCHLEAR IMPLANT CONFERENCE  ·  ARABCIC 2026", "size": 14,
      "bold": True, "color": GOLD, "align": PP_ALIGN.CENTER}])
txt(s, 0.8, 1.28, 11.73, 0.35,
    [{"text": "Jeddah, Kingdom of Saudi Arabia  ·  26–28 November 2026", "size": 12,
      "color": RGBColor(0xBF, 0xD0, 0xE0), "align": PP_ALIGN.CENTER}])
txt(s, 0.8, 2.05, 11.73, 1.75,
    [{"text": "A Fuzzy Logic Framework for Sensorineural Hearing Loss Classification",
      "size": 33, "bold": True, "color": WHITE, "font": HDR, "align": PP_ALIGN.CENTER}])
txt(s, 0.8, 3.62, 11.73, 0.5,
    [{"text": "Preserving diagnostic gradation lost to crisp WHO thresholds",
      "size": 17, "italic": True, "color": GOLD, "align": PP_ALIGN.CENTER}])
txt(s, 0.8, 4.55, 11.73, 0.85,
    [{"text": "Dr. Sanyaolu Ameye (MBBS, FWACS, FMCORL, Pg Cert ML/AI, MSc Data Science)",
      "size": 15, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER},
     {"text": "ENT Surgeon and Data Scientist", "size": 13, "color": RGBColor(0xBF, 0xD0, 0xE0),
      "align": PP_ALIGN.CENTER, "space_after": 2},
     {"text": "King Fahad Specialist Hospital, Tabuk, Kingdom of Saudi Arabia",
      "size": 13, "color": RGBColor(0xBF, 0xD0, 0xE0), "align": PP_ALIGN.CENTER}])
chip(s, 2.20, 5.80, 2.6, 1.05, "10,889", "adults 20–69 y · 3 NHANES cycles", big_size=24)
chip(s, 5.37, 5.80, 2.6, 1.05, "48", "audiology-derived fuzzy rules", big_size=24)
chip(s, 8.54, 5.80, 2.6, 1.05, "0–100", "continuous FAI scale", big_size=24)
notes(s, "Good morning. Hearing loss classification has barely changed since Goodman in 1965. "
          "This talk presents a fuzzy logic framework that preserves the diagnostic gradation "
          "that crisp WHO thresholds discard, validated on three NHANES adult cycles (20–69 y).")

# ================================================================ SLIDE 2 — PROBLEM (light)
s = add_slide(LIGHT)
title_bar(s, "The problem: crisp thresholds on continuous biology", 2)
# left card
rect(s, 0.55, 1.25, 6.0, 5.55, CARD, line=BORDER, rounded=True, radius=0.045)
txt(s, 0.80, 1.45, 5.5, 0.4,
    [{"text": "One decibel changes the label", "size": 17, "bold": True, "color": NAVY, "font": HDR}])
# patient A
pa = rect(s, 0.80, 1.95, 2.65, 1.55, SOFT_T, rounded=True, radius=0.08)
txt(s, 0.95, 2.05, 2.35, 1.4,
    [{"text": "PATIENT A", "size": 11, "bold": True, "color": TEAL},
     {"text": "PTA-4 = 25 dB HL", "size": 15, "bold": True, "color": INK, "space_before": 4},
     {"text": "Label: Normal", "size": 13, "color": INK, "space_before": 3}])
pb = rect(s, 3.70, 1.95, 2.65, 1.55, SOFT_G, rounded=True, radius=0.08)
txt(s, 3.85, 2.05, 2.35, 1.4,
    [{"text": "PATIENT B", "size": 11, "bold": True, "color": RGBColor(0xB8, 0x8A, 0x2E)},
     {"text": "PTA-4 = 26 dB HL", "size": 15, "bold": True, "color": INK, "space_before": 4},
     {"text": "Label: Mild hearing loss", "size": 13, "color": INK, "space_before": 3}])
txt(s, 0.80, 3.65, 5.5, 0.75,
    [{"text": "One decibel, well within the test–retest variability of audiometry (±5–10 dB), "
               "produces categorically different labels, follow-up and candidacy decisions.",
      "size": 12.5, "color": INK, "line_spacing": 1.15}])
rect(s, 0.80, 4.50, 5.55, 2.1, RGBColor(0xEF, 0xF4, 0xF8), line=BORDER, rounded=True, radius=0.07)
txt(s, 0.98, 4.62, 5.2, 1.9,
    [{"text": "What a single PTA number discards", "size": 13, "bold": True, "color": NAVY},
     {"text": "•  A 4 kHz notch and a flat 30 dB loss share one grade, with very different "
               "functional profiles", "size": 12, "color": INK, "space_before": 5, "line_spacing": 1.12},
     {"text": "•  Configuration and asymmetry are invisible to a pure-tone average",
      "size": 12, "color": INK, "space_before": 5, "line_spacing": 1.12}])
# right stats
c1 = rect(s, 6.85, 1.25, 5.95, 2.35, CARD, line=BORDER, rounded=True, radius=0.06)
txt(s, 7.05, 1.45, 5.55, 1.0,
    [{"text": "39.2%", "size": 44, "bold": True, "color": NAVY, "font": HDR}])
txt(s, 7.05, 2.45, 5.55, 1.05,
    [{"text": "of ears fall within ±5 dB of a WHO severity boundary — the zone where "
               "crisp labels are most fragile", "size": 13, "color": INK, "line_spacing": 1.15}])
c2 = rect(s, 6.85, 3.80, 5.95, 1.75, CARD, line=BORDER, rounded=True, radius=0.06)
txt(s, 7.05, 4.00, 5.55, 0.8,
    [{"text": "±5–10 dB", "size": 30, "bold": True, "color": TEAL, "font": HDR}])
txt(s, 7.05, 4.80, 5.55, 0.7,
    [{"text": "test–retest variability of pure-tone audiometry — smaller than the gaps "
               "crisp rules insist on", "size": 12, "color": INK, "line_spacing": 1.12}])
txt(s, 6.85, 5.85, 5.95, 0.9,
    [{"text": "Goodman (1965) → Clark (1981) → WHO World Report on Hearing (2021): "
               "the classification has barely changed in fifty years.",
      "size": 11.5, "italic": True, "color": MUTED, "line_spacing": 1.15}])
notes(s, "Two patients separated by a single decibel are categorically different under WHO "
          "rules, even though a decibel sits inside audiometry's own test-retest error. "
          "Nearly four in ten ears in NHANES sit within 5 dB of a boundary. And the PTA "
          "throws away frequency-specific detail and configuration entirely.")

# ================================================================ SLIDE 3 — SOLUTION (light)
s = add_slide(LIGHT)
title_bar(s, "The solution: fuzzy sets for graded severity", 3)
# crisp card
rect(s, 0.55, 1.30, 6.0, 4.5, CARD, line=BORDER, rounded=True, radius=0.05)
card_header(s, 0.55, 1.30, 6.0, 0.55, "CRISP LOGIC — TODAY", fill=NAVY)
txt(s, 0.85, 2.05, 5.45, 3.6,
    [{"runs": [{"text": "✗  ", "size": 14, "bold": True, "color": RGBColor(0xC0, 0x39, 0x2B)},
               {"text": "Binary membership: a threshold is in one category or not at all",
                "size": 13.5, "color": INK}], "space_after": 11, "line_spacing": 1.12},
     {"runs": [{"text": "✗  ", "size": 14, "bold": True, "color": RGBColor(0xC0, 0x39, 0x2B)},
               {"text": "Sharp cutoffs at 25 / 40 / 55 / 70 / 90 dB", "size": 13.5, "color": INK}],
      "space_after": 11},
     {"runs": [{"text": "✗  ", "size": 14, "bold": True, "color": RGBColor(0xC0, 0x39, 0x2B)},
               {"text": "No frequency resolution beyond an average", "size": 13.5, "color": INK}],
      "space_after": 11},
     {"runs": [{"text": "✗  ", "size": 14, "bold": True, "color": RGBColor(0xC0, 0x39, 0x2B)},
               {"text": "No configuration typing or asymmetry scoring", "size": 13.5, "color": INK}],
      "space_after": 0}])
# fuzzy card
rect(s, 6.80, 1.30, 6.0, 4.5, CARD, line=BORDER, rounded=True, radius=0.05)
card_header(s, 6.80, 1.30, 6.0, 0.55, "FUZZY LOGIC — THIS STUDY", fill=TEAL)
txt(s, 7.10, 2.05, 5.45, 3.6,
    [{"runs": [{"text": "✓  ", "size": 14, "bold": True, "color": TEAL},
               {"text": "Gradual membership (0 → 1) in every severity category",
                "size": 13.5, "color": INK}], "space_after": 11, "line_spacing": 1.12},
     {"runs": [{"text": "✓  ", "size": 14, "bold": True, "color": TEAL},
               {"text": "Overlapping trapezoidal functions at each boundary (2.0 dB overlap)",
                "size": 13.5, "color": INK}], "space_after": 11},
     {"runs": [{"text": "✓  ", "size": 14, "bold": True, "color": TEAL},
               {"text": "Continuous FAI score (0–100), frequency-resolved", "size": 13.5, "color": INK}],
      "space_after": 11},
     {"runs": [{"text": "✓  ", "size": 14, "bold": True, "color": TEAL},
               {"text": "Configuration vector + asymmetry index", "size": 13.5, "color": INK}],
      "space_after": 0}])
rect(s, 0.55, 6.00, 12.25, 0.85, TEAL, rounded=True, radius=0.10)
txt(s, 0.75, 6.00, 11.85, 0.85,
    [{"text": "Fuzzy Audiometric Index (FAI): continuous 0–100 · frequency-resolved · "
               "interpretable rule base", "size": 14.5, "bold": True, "color": WHITE,
      "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
notes(s, "Fuzzy set theory, introduced by Zadeh, lets a threshold belong to several severity "
          "categories at once. The result is a continuous Fuzzy Audiometric Index, a "
          "configuration vector, and an asymmetry index — all produced by an interpretable "
          "rule base rather than a black box.")

# ================================================================ SLIDE 4 — METHODS (light)
s = add_slide(LIGHT)
title_bar(s, "Methods: data and the Mamdani fuzzy inference system", 4)
rect(s, 0.55, 1.30, 5.6, 5.5, CARD, line=BORDER, rounded=True, radius=0.05)
txt(s, 0.80, 1.50, 5.1, 0.4,
    [{"text": "Data: 3 NHANES adult cycles", "size": 16, "bold": True, "color": NAVY, "font": HDR}])
txt(s, 0.80, 1.95, 5.1, 4.7,
    [{"text": "•  Adults 20–69 y: 3 NHANES cycles (1999-2000, 2011-12, 2015-16), n = 10,889", "size": 12.5,
      "color": INK, "space_after": 8, "line_spacing": 1.12},
     {"text": "•  80/20 participant-level split: 15,656 training / 3,912 held-out test ears (no bilateral leakage)",
      "size": 12.5, "color": INK, "space_after": 8, "line_spacing": 1.12},
     {"text": "•  Membership functions and rule base developed on the training set only",
      "size": 12.5, "color": INK, "space_after": 8, "line_spacing": 1.12},
     {"text": "•  All metrics reported on the held-out test set", "size": 12.5, "color": INK,
      "space_after": 10, "line_spacing": 1.12},
     {"text": "Comparators", "size": 13.5, "bold": True, "color": TEAL, "space_after": 5},
     {"text": "•  PTA-4 (WHO) crisp classification", "size": 12.5, "color": INK, "space_after": 5},
     {"text": "•  XGBoost and Random Forest", "size": 12.5, "color": INK, "space_after": 10},
     {"text": "Implementation: scikit-fuzzy (v0.5.0), scikit-learn, XGBoost, SciPy",
      "size": 11.5, "italic": True, "color": MUTED, "line_spacing": 1.1}])
img, iw, ih = pic(s, f"{FIG}/fig2_fis_architecture.png", 6.45, 1.60, w=6.35)
txt(s, 6.45, 1.60 + ih + 0.12, 6.35, 0.7,
    [{"text": "Mamdani FIS: 7 frequency thresholds → fuzzification → 48 rules → "
               "centroid defuzzification → FAI", "size": 11.5, "italic": True, "color": MUTED,
      "align": PP_ALIGN.CENTER, "line_spacing": 1.1}])
notes(s, "We combined three NHANES cycles that tested adults 20-69 y (10,889 participants; 19,568 clean ears). "
          "An 80/20 stratified split keeps every WHO category in both partitions, and "
          "everything is optimised on training and reported on the held-out test set. "
          "Comparators are the WHO PTA-4 rules, XGBoost and Random Forest.")

# ================================================================ SLIDE 5 — FUZZIFICATION (light)
s = add_slide(LIGHT)
title_bar(s, "Fuzzification: trapezoidal membership functions", 5)
img, iw, ih = pic(s, f"{FIG}/fig1_membership_functions.png", 0.55, 1.35, w=6.5)
txt(s, 0.55, 1.35 + ih + 0.10, 6.5, 0.6,
    [{"text": "Overlapping trapezoids for six severity categories, with NHANES density overlay",
      "size": 11, "italic": True, "color": MUTED, "align": PP_ALIGN.CENTER}])
rect(s, 0.55, 5.30, 6.5, 1.55, SOFT_T, line=BORDER, rounded=True, radius=0.07)
txt(s, 0.75, 5.45, 6.1, 1.3,
    [{"text": "Why trapezoids", "size": 12.5, "bold": True, "color": TEAL},
     {"text": "NHANES severity distributions are asymmetric with plateau cores — "
               "poorly captured by triangular or Gaussian shapes.", "size": 11.5,
      "color": INK, "space_before": 4, "line_spacing": 1.12}])
# param table
data = [["Category", "a", "b", "c", "d"],
        ["Normal", "0.0", "6.4", "16.4", "27.8"],
        ["Mild", "25.8", "26.3", "36.4", "43.5"],
        ["Moderate", "41.5", "42.0", "54.3", "58.5"],
        ["Moderately Severe", "56.5", "57.0", "67.9", "73.5"],
        ["Severe", "71.5", "72.0", "83.7", "93.0"],
        ["Profound", "91.0", "91.5", "104.6", "120.0"]]
make_table(s, 7.30, 1.35, 5.5, data, [2.35, 0.79, 0.79, 0.79, 0.79], row_h=0.42,
           body_size=10.5, header_size=11)
txt(s, 7.30, 4.50, 5.5, 0.35,
    [{"text": "Optimised on the NHANES training distribution [a, b, c, d]",
      "size": 10, "italic": True, "color": MUTED}])
rect(s, 7.30, 4.95, 5.5, 1.9, CARD, line=BORDER, rounded=True, radius=0.07)
txt(s, 7.50, 5.10, 5.1, 1.6,
    [{"text": "Exactly 2.0 dB overlap at each WHO boundary → smooth transitions",
      "size": 12, "color": INK, "space_after": 8, "line_spacing": 1.12},
     {"runs": [{"text": "26 dB HL → ", "size": 12, "color": INK},
               {"text": "Normal μ 0.16", "size": 12, "bold": True, "color": TEAL},
               {"text": " and ", "size": 12, "color": INK},
               {"text": "Mild μ 0.40", "size": 12, "bold": True, "color": TEAL},
               {"text": ": borderline, not binary", "size": 12, "color": INK}],
      "line_spacing": 1.12}])
notes(s, "Each frequency band has six overlapping trapezoidal membership functions. "
          "The optimised parameters sit a few dB to the right of the classic WHO cutoffs, "
          "with exactly 2 dB of overlap, so a 26 dB threshold is simultaneously 16% normal "
          "and 40% mild.")

# ================================================================ SLIDE 6 — RULE BASE (light)
s = add_slide(LIGHT)
title_bar(s, "Rule base: 48 audiology-derived rules in four groups", 6)
groups = [
    ("SEVERITY · 12 RULES", "Map thresholds to severity output; blended rules handle "
     "transitional zones.", "IF threshold is Mild AND slope is Flat\nTHEN severity is Mild", NAVY),
    ("CONFIGURATION · 14 RULES", "Combine slope and notch to type the audiogram shape.",
     "IF slope is Steeply Sloping AND notch is Deep\nTHEN configuration is Notched", TEAL),
    ("ASYMMETRY · 12 RULES", "Score inter-aural differences and adjust severity.",
     "IF inter-aural difference > 30 dB at 3+ frequencies\nTHEN asymmetry is Severely Asymmetric", NAVY),
    ("MIXED-LOSS · 10 RULES", "Handle complex presentations combining all dimensions.",
     "IF asymmetric AND notched\nTHEN flag mixed aetiology", TEAL),
]
pos = [(0.55, 1.35), (6.85, 1.35), (0.55, 3.85), (6.85, 3.85)]
for (label, desc, ex, col), (x, y) in zip(groups, pos):
    rect(s, x, y, 5.95, 2.30, CARD, line=BORDER, rounded=True, radius=0.06)
    txt(s, x + 0.25, y + 0.18, 5.45, 0.4,
        [{"text": label, "size": 13.5, "bold": True, "color": col}])
    txt(s, x + 0.25, y + 0.60, 5.45, 0.75, [{"text": desc, "size": 11.5, "color": INK,
        "line_spacing": 1.12}])
    txt(s, x + 0.25, y + 1.42, 5.45, 0.75,
        [{"text": ex, "size": 10.5, "italic": True, "color": MUTED, "line_spacing": 1.15}])
rect(s, 0.55, 6.35, 12.25, 0.65, NAVY, rounded=True, radius=0.12)
txt(s, 0.75, 6.35, 11.85, 0.65,
    [{"text": "Mamdani min–max composition · centroid-of-area defuzzification · FAI 0–100",
      "size": 13, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}],
    anchor=MSO_ANCHOR.MIDDLE)
notes(s, "Forty-eight expert-derived rules in four groups: severity, configuration, "
          "asymmetry and mixed-loss patterns. Inference uses Mamdani min-max composition "
          "with centroid-of-area defuzzification to produce the continuous FAI.")

# ================================================================ SLIDE 7 — RESULTS (light)
s = add_slide(LIGHT)
title_bar(s, "Results: FAI vs WHO PTA-4 and ML comparators", 7)
data = [["Method", "Weighted κ", "Borderline (±5 dB)", "Clear-case"],
        ["FAI (fuzzy)", "0.93", "79.8%", "98.1%"],
        ["PTA-4 (WHO)", "1.00*", "100%*", "100%*"],
        ["XGBoost", "0.98", "93.4%", "100%"],
        ["Random Forest", "0.98", "93.6%", "100%"]]
make_table(s, 0.55, 1.35, 7.6, data, [2.5, 1.7, 1.7, 1.7], row_h=0.5,
           body_size=12, header_size=12, highlight_rows={2})
txt(s, 0.55, 3.85, 7.6, 0.35,
    [{"text": "*PTA-4 is the reference standard; borderline accuracy reflects discrepancy "
               "within ±5 dB of thresholds.", "size": 9.5, "italic": True, "color": MUTED}])
rect(s, 0.55, 4.30, 7.6, 1.35, SOFT_T, line=BORDER, rounded=True, radius=0.07)
txt(s, 0.75, 4.45, 7.2, 1.1,
    [{"text": "In borderline cases (±5 dB), the fuzzy system deliberately reclassifies "
               "(79.8% agree with the crisp grade) while matching it in clear cases (98.1%).",
      "size": 12.5, "color": INK, "line_spacing": 1.2}])
img, iw, ih = pic(s, f"{FIG}/fig3_bland_altman.png", 8.45, 1.35, w=4.35)
txt(s, 8.45, 1.35 + ih + 0.08, 4.35, 0.35,
    [{"text": "Bland–Altman: FAI vs PTA-4 reference", "size": 10, "italic": True,
      "color": MUTED, "align": PP_ALIGN.CENTER}])
chip(s, 8.45, 5.05, 4.35, 0.95, "κ = 0.93", "vs WHO PTA-4 reference, 20–69 y test", big_size=20,
     small_size=10, fill=NAVY, big_color=GOLD)
chip(s, 8.45, 6.10, 4.35, 0.85, "ρ = 0.81 · MAE 5.0 dB", "FAI vs PTA-4, n = 3,912 test ears",
     big_size=15, small_size=10, fill=NAVY2, big_color=WHITE)
notes(s, "On the combined 20-69 y test set the FAI shows substantial agreement with the "
          "WHO PTA-4 reference (kappa 0.93, Spearman 0.80) and matches the crisp grade "
          "on 98.1% of clear cases. 18.7% of test ears sit within 5 dB of a severity "
          "boundary; there the system returns graded membership (79.8% agreement with "
          "the crisp label) rather than forcing a binary choice. Spearman is attenuated "
          "by range restriction (88.6% normal). The ML comparators regress PTA-4 itself "
          "and are reference-in-disguise, not independent benchmarks.")

# ================================================================ SLIDE 8 — 25/26 dB (light)
s = add_slide(LIGHT)
title_bar(s, "The 25/26 dB boundary: graded, not binary", 8)
data = [["PTA-4", "Crisp label", "Fuzzy label", "FAI", "Normal μ", "Mild μ"],
        ["24 dB", "Normal", "Normal", "11.3", "0.33", "0.00"],
        ["25 dB", "Normal", "Normal", "11.6", "0.25", "0.00"],
        ["26 dB", "Mild", "Mild", "25.7", "0.16", "0.40"],
        ["30 dB", "Mild", "Mild", "30.0", "0.00", "1.00"]]
make_table(s, 0.55, 1.35, 6.4, data, [1.05, 1.05, 1.05, 0.85, 1.2, 1.2], row_h=0.46,
           body_size=11, header_size=11, highlight_rows={3})
rect(s, 0.55, 4.60, 6.4, 2.05, CARD, line=BORDER, rounded=True, radius=0.06)
txt(s, 0.75, 4.75, 6.0, 1.8,
    [{"text": "A patient with PTA 26 dB is simultaneously 16% normal and 40% mild — "
               "the crisp label hides this ambiguity.", "size": 12.5, "color": INK,
      "space_after": 8, "line_spacing": 1.2},
     {"text": "FAI rises smoothly from 11.3 at 24 dB to 30.0 at 30 dB; crisp rules impose a step at 26 dB.",
      "size": 12.5, "color": INK, "line_spacing": 1.2}])
img, iw, ih = pic(s, f"{FIG}/fig5_borderline_analysis.png", 7.25, 1.35, w=5.55)
txt(s, 7.25, 1.35 + ih + 0.08, 5.55, 0.6,
    [{"text": "The fuzzy system deviates from crisp labels in the borderline zone "
               "(agreement 80–84% within ±1–5 dB) while matching crisp in clear cases (98.1%)", "size": 10.5, "italic": True,
      "color": MUTED, "align": PP_ALIGN.CENTER, "line_spacing": 1.1}])
notes(s, "This is the clinical core. At the normal-mild boundary the fuzzy system reports "
          "graded membership instead of a forced binary label; it matches the crisp label "
          "in clear cases (98.1%) while deliberately deviating from it in the borderline "
          "zone (agreement 83–86% within ±1–3 dB).")

# ================================================================ SLIDE 9 — CASES (light)
s = add_slide(LIGHT)
title_bar(s, "Clinical case studies from NHANES", 9)
img, iw, ih = pic(s, f"{FIG}/fig4_clinical_cases.png", 0.55, 1.30, w=6.05)
txt(s, 0.55, 1.30 + ih + 0.06, 6.05, 0.3,
    [{"text": "Four representative audiograms with membership bar charts", "size": 9.5,
      "italic": True, "color": MUTED, "align": PP_ALIGN.CENTER}])
cases = [
    ("A · BORDERLINE", "PTA 27.5 dB “Mild” · FAI 29.4", "Normal μ 0.04, Mild μ 1.00 — graded, not forced"),
    ("B · MASKED NOTCH", "PTA 23 dB “Normal”", "25 dB notch at 4 kHz · Notched μ 0.78 — PTA missed it"),
    ("C · PRESBYCUSIS", "PTA 33.8 dB “Mild” · FAI 30.0", "Sloping 35 dB — frequency gradient preserved"),
    ("D · ASYMMETRY", "R 53.8 dB Mod / L 28.8 Mild", "Composite FAI 54.5 (Mod-Sev) — impact quantified"),
]
cpos = [(6.95, 1.30), (6.95, 3.10), (6.95, 4.90), (6.95, 6.05)]
for (lab, head, sub), (x, y) in zip(cases, cpos):
    hgt = 1.65 if lab != "D · ASYMMETRY" else 1.15
    rect(s, x, y, 5.85, hgt, CARD, line=BORDER, rounded=True, radius=0.08)
    txt(s, x + 0.2, y + 0.1, 5.45, hgt - 0.15,
        [{"text": lab, "size": 11, "bold": True, "color": TEAL},
         {"text": head, "size": 12, "bold": True, "color": INK, "space_before": 3},
         {"text": sub, "size": 10.5, "color": MUTED, "space_before": 3, "line_spacing": 1.1}])
notes(s, "Four cases illustrate the clinical value. A borderline patient gets graded "
          "membership rather than a forced label. A noise notch is caught despite a "
          "normal PTA. Presbycusis keeps its frequency gradient. And asymmetric loss "
          "has its impact quantified and folded into the severity estimate.")

# ================================================================ SLIDE 10 — CONFIGURATION (light)
s = add_slide(LIGHT)
title_bar(s, "Configuration typing and asymmetry", 10)
img, iw, ih = pic(s, f"{FIG}/fig_eda_config_dist.png", 0.55, 1.35, w=4.9)
txt(s, 0.55, 1.35 + ih + 0.08, 4.9, 0.6,
    [{"text": "Configuration distribution, combined adult cohort (slope = 4 kHz − 500 Hz)", "size": 10.5, "italic": True, "color": MUTED,
      "align": PP_ALIGN.CENTER, "line_spacing": 1.1}])
rows = [
    ("Flat 56.7% · Gently sloping 18.0%", "dominant configurations (4 kHz − 500 Hz)", NAVY, GOLD),
    ("Steeply 8.8% · Precipitous 3.5% · Rising 13.0%", "exploratory output; accuracy table not reproducible", TEAL, WHITE),
    ("6-shape membership vector", "graded configuration typing per ear", NAVY2, WHITE),
    ("5.1% of participants", "exceed 15 dB inter-aural difference; continuous index feeds severity "
     "adjustment (Case D)", TEAL, WHITE),
]
y = 1.35
for big, small, fill, bcol in rows:
    chip(s, 5.75, y, 7.05, 1.15, big, small, fill=fill, big_color=bcol, small_color=WHITE,
         big_size=19, small_size=10.5)
    y += 1.32
notes(s, "Configuration typing is an exploratory output: the combined adult cohort is "
          "predominantly flat and gently sloping, and the fuzzy shape rules emit a "
          "six-category membership vector. Asymmetry affects 5.1% of the cohort at the "
          "15 dB threshold, and the fuzzy index quantifies it continuously.")

# ================================================================ SLIDE 11 — LONGITUDINAL (light)
s = add_slide(LIGHT)
title_bar(s, "Longitudinal tracking: sub-threshold change", 11)
img, iw, ih = pic(s, f"{FIG}/fig7_longitudinal.png", 0.55, 1.30, w=12.25)
txt(s, 0.55, 1.30 + ih + 0.08, 12.25, 0.35,
    [{"text": "FAI trajectories (coloured) vs PTA grade (grey steps) over 4–8 years — "
               "smooth progression instead of discrete jumps", "size": 10.5, "italic": True,
      "color": MUTED, "align": PP_ALIGN.CENTER}])
stats = [
    ("+1.7 pts/yr", "mean annual FAI change, age-related subgroup"),
    ("1 grade / 4.2 yr", "PTA steps between categories"),
    ("28%", "FAI shifts ≥ 5 pts without crossing a PTA boundary"),
]
x = 0.55
for big, small in stats:
    chip(s, x, 5.55, 3.95, 1.25, big, small, fill=NAVY, big_color=GOLD, big_size=20, small_size=10.5)
    x += 4.18
notes(s, "In the 120-participant longitudinal sub-cohort the FAI moves continuously, "
          "about 1.7 points per year in age-related loss, while PTA grades step once "
          "every 4.2 years. 28% of participants showed a change of 5 or more FAI points "
          "without crossing any PTA boundary — sub-threshold shifts that include both "
          "ototoxicity and recovery after acute noise exposure.")

# ================================================================ SLIDE 12 — IMPLICATIONS (light)
s = add_slide(LIGHT)
title_bar(s, "Clinical implications", 12)
imp = [
    ("Audiologic triage", "Flag patients straddling severity boundaries for monitoring "
     "instead of forcing one label"),
    ("Hearing aid fitting", "Continuous FAI maps to gain; configuration vector guides "
     "frequency-specific amplification"),
    ("Ototoxicity monitoring", "≥ 5-point FAI shifts (cisplatin, aminoglycosides) before "
     "PTA crosses a boundary"),
    ("Noise surveillance", "Notch detection flags early noise damage despite a normal PTA "
     "(Case B)"),
    ("Low-resource settings", "Runs on mobile audiometry (hearWHO, Mimi); interpretable "
     "and computationally light"),
    ("Research endpoints", "Continuous outcome for hearing-preservation and "
     "rehabilitation trials"),
]
pos = [(0.55, 1.35), (4.72, 1.35), (8.89, 1.35), (0.55, 4.05), (4.72, 4.05), (8.89, 4.05)]
for (head, body), (x, y) in zip(imp, pos):
    rect(s, x, y, 3.92, 2.45, CARD, line=BORDER, rounded=True, radius=0.07)
    rect(s, x, y, 0.09, 2.45, TEAL)
    txt(s, x + 0.25, y + 0.18, 3.5, 0.4, [{"text": head, "size": 13.5, "bold": True, "color": NAVY}])
    txt(s, x + 0.25, y + 0.62, 3.5, 1.7, [{"text": body, "size": 11.5, "color": INK, "line_spacing": 1.15}])
notes(s, "The framework has practical routes into care: more honest triage, gain "
          "programming for hearing aids, earlier ototoxicity detection, noise surveillance, "
          "low-resource mobile screening, and a continuous endpoint for trials.")

# ================================================================ SLIDE 13 — LIMITATIONS (light)
s = add_slide(LIGHT)
title_bar(s, "Limitations and next steps", 13)
rect(s, 0.55, 1.30, 6.0, 5.5, CARD, line=BORDER, rounded=True, radius=0.05)
txt(s, 0.80, 1.50, 5.5, 0.4, [{"text": "Limitations", "size": 16, "bold": True, "color": NAVY, "font": HDR}])
txt(s, 0.80, 1.95, 5.5, 4.7,
    [{"text": "•  Air-conduction thresholds only — bone conduction needed for mixed-loss detection",
      "size": 12, "color": INK, "space_after": 7, "line_spacing": 1.12},
     {"text": "•  Rule base expert-derived, not yet through formal Delphi consensus",
      "size": 12, "color": INK, "space_after": 7, "line_spacing": 1.12},
     {"text": "•  NHANES skews toward mild–moderate age-related loss", "size": 12, "color": INK,
      "space_after": 7, "line_spacing": 1.12},
     {"text": "•  No extended high-frequency audiometry (10–16 kHz)", "size": 12, "color": INK,
      "space_after": 7, "line_spacing": 1.12},
     {"text": "•  Reference standard (WHO PTA-4) carries the same boundary issue",
      "size": 12, "color": INK, "space_after": 7, "line_spacing": 1.12},
     {"text": "•  Outcome validation pending (hearing aids, CI candidacy)", "size": 12, "color": INK,
      "space_after": 0, "line_spacing": 1.12}])
rect(s, 6.85, 1.30, 6.0, 5.5, CARD, line=BORDER, rounded=True, radius=0.05)
txt(s, 7.10, 1.50, 5.5, 0.4, [{"text": "Next steps", "size": 16, "bold": True, "color": TEAL, "font": HDR}])
txt(s, 7.10, 1.95, 5.5, 4.7,
    [{"text": "•  Multi-site validation — paediatric and Menière’s cohorts", "size": 12,
      "color": INK, "space_after": 7, "line_spacing": 1.12},
     {"text": "•  Mobile audiometry integration for community screening in low-resource settings",
      "size": 12, "color": INK, "space_after": 7, "line_spacing": 1.12},
     {"text": "•  Web-based FAI calculator with REST API for EHR integration", "size": 12,
      "color": INK, "space_after": 7, "line_spacing": 1.12},
     {"text": "•  Neuro-fuzzy learning (ANFIS) — learn rule parameters from data while staying "
               "interpretable", "size": 12, "color": INK, "space_after": 7, "line_spacing": 1.12},
     {"text": "•  FAI as a continuous feature for ML outcome models (speech-in-noise, hearing aid "
               "outcomes, CI candidacy)", "size": 12, "color": INK, "space_after": 0, "line_spacing": 1.12}])
notes(s, "Honest limitations: air-conduction only, no Delphi consensus yet, NHANES's "
          "mild-to-moderate skew, no extended high frequencies, and outcome validation "
          "still pending. Next: multi-site validation, mobile integration, a web "
          "calculator, ANFIS, and using the FAI as a feature in outcome models.")

# ================================================================ SLIDE 14 — CONCLUSION (dark)
s = add_slide(NAVY)
waveform(s, 0, 7.28, 13.333, color=GOLD, max_h=0.5)
txt(s, 0.8, 0.85, 11.73, 0.4,
    [{"text": "TAKE-HOME", "size": 14, "bold": True, "color": GOLD, "align": PP_ALIGN.CENTER}])
concl = [
    "A continuous, frequency-resolved Fuzzy Audiometric Index preserves the gradation "
    "that crisp PTA thresholds discard.",
    "κ = 0.93 against WHO PTA-4; 98.1% agreement in clear cases, with deliberate "
    "graded reclassification in the borderline zone (79.8%).",
    "Interpretable by design — a rule base clinicians can inspect, question and modify.",
]
y = 1.75
for i, c in enumerate(concl, 1):
    rect(s, 1.30, y, 10.73, 1.15, NAVY2, rounded=True, radius=0.10)
    txt(s, 1.60, y, 1.0, 1.15, [{"text": str(i), "size": 28, "bold": True, "color": GOLD,
        "font": HDR, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 2.55, y, 9.2, 1.15, [{"text": c, "size": 14.5, "color": WHITE, "line_spacing": 1.18}],
        anchor=MSO_ANCHOR.MIDDLE)
    y += 1.32
txt(s, 1.30, 5.95, 10.73, 0.9,
    [{"text": "“As audiology moves toward personalised, data-driven care, fuzzy logic keeps "
               "the continuity of auditory function front and centre.”",
      "size": 14.5, "italic": True, "color": GOLD, "align": PP_ALIGN.CENTER, "line_spacing": 1.2}])
notes(s, "Three messages: the FAI preserves gradation, it deliberately reclassifies "
          "borderline cases where crisp labels force a binary choice, and it stays "
          "interpretable throughout.")

# ================================================================ SLIDE 15 — THANK YOU (dark)
s = add_slide(NAVY)
waveform(s, 0, 7.28, 13.333, color=GOLD, max_h=0.55)
waveform(s, 0, 7.28, 13.333, color=TEAL, max_h=0.3)
txt(s, 0.8, 1.9, 11.73, 1.0,
    [{"text": "Thank you", "size": 48, "bold": True, "color": WHITE, "font": HDR,
      "align": PP_ALIGN.CENTER}])
txt(s, 0.8, 2.95, 11.73, 0.5,
    [{"text": "Questions and discussion", "size": 18, "color": RGBColor(0xBF, 0xD0, 0xE0),
      "align": PP_ALIGN.CENTER}])
txt(s, 0.8, 4.15, 11.73, 0.9,
    [{"text": "Dr. Sanyaolu Ameye (MBBS, FWACS, FMCORL, Pg Cert ML/AI, MSc Data Science)",
      "size": 15, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER},
     {"text": "ENT Surgeon and Data Scientist", "size": 13, "color": RGBColor(0xBF, 0xD0, 0xE0),
      "align": PP_ALIGN.CENTER, "space_after": 3},
     {"text": "King Fahad Specialist Hospital, Tabuk, Kingdom of Saudi Arabia",
      "size": 13, "color": RGBColor(0xBF, 0xD0, 0xE0), "align": PP_ALIGN.CENTER}])
txt(s, 0.8, 5.25, 11.73, 0.45,
    [{"text": "sanyaameye@hotmail.com", "size": 15, "bold": True, "color": GOLD,
      "align": PP_ALIGN.CENTER}])
txt(s, 0.8, 5.85, 11.73, 0.4,
    [{"text": "Code and data: github.com/ameye/fuzzy-audiogram", "size": 12,
      "color": RGBColor(0x9A, 0xB0, 0xC6), "align": PP_ALIGN.CENTER}])
notes(s, "Thank you. Happy to take questions on the rule base, the validation, or "
          "deployment plans for cochlear implant and hearing aid pathways.")

# ---------------------------------------------------------------- save
outdir = "/opt/data/fuzzy-audiogram/arabcic2026"
os.makedirs(outdir, exist_ok=True)
out = f"{outdir}/arabcic2026_fuzzy_audiogram.pptx"
prs.save(out)
print("saved:", out, os.path.getsize(out), "bytes")
